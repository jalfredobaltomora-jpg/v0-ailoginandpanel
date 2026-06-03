#!/usr/bin/env python3
"""
Export Engine Monitor - Executive PowerPoint
==============================================
Uses python-pptx to generate a 4-slide presentation with:
- Dark theme (#0b111b background, #00d084 accents)
- Slide 1: Title + system logo
- Slide 2: 4 KPI indicator cards with colored icons
- Slide 3: Charts as high-res images (bar + donut)
- Slide 4: Detail table with dark alternating rows

Usage:
  python scripts/export_engine_monitor_ppt.py data.json chart.png -o report.pptx
  python scripts/export_engine_monitor_ppt.py --sample -o report.pptx
"""

import json, sys, argparse, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
BG_DARK = RGBColor(0x0B, 0x11, 0x1B)
GREEN_NEON = RGBColor(0x00, 0xD0, 0x84)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x8B, 0x94, 0x9E)
DARK_CARD = RGBColor(0x16, 0x1B, 0x22)
RED_ACCENT = RGBColor(0xF4, 0x43, 0x36)
BLUE_ACCENT = RGBColor(0x06, 0xB6, 0xD4)
AMBER_ACCENT = RGBColor(0xF5, 0x9E, 0x0B)
EMERALD_ACCENT = RGBColor(0x10, 0xB9, 0x81)


def add_dark_bg(slide):
    """Fill slide background with dark color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_card(slide, left, top, width, height, title, value, color=GREEN_NEON):
    """Add a KPI card with dark background and colored accent."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_CARD
    shape.line.color.rgb = color
    shape.line.width = Pt(2)

    # Value
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()

    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(28)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY_TEXT
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER


def build_presentation(records, chart_image_path, output_path, filters=None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Compute aggregates
    total_audits = sum(r.get("audits", 0) for r in records)
    total_failures = sum(r.get("failures", 0) for r in records)
    total_meas_def = sum(r.get("measDef", 0) for r in records)
    total_vis_def = sum(r.get("visDef", 0) for r in records)
    total_meas_ins = sum(r.get("measIns", 0) for r in records)
    total_vis_ins = sum(r.get("visIns", 0) for r in records)

    fr_pct = (total_failures / total_audits * 100) if total_audits > 0 else 0
    mr_pct = (total_meas_def / total_meas_ins * 100) if total_meas_ins > 0 else 0
    vr_pct = (total_vis_def / total_vis_ins * 100) if total_vis_ins > 0 else 0

    def kpi_color(pct):
        if pct < 5: return GREEN_NEON
        if pct < 15: return AMBER_ACCENT
        return RED_ACCENT

    # Build filter header text
    filter_header = ""
    if filters:
        parts = []
        for key, label in [("months", "Meses"), ("factories", "Fabricas"),
                           ("buyers", "Clientes"), ("points", "Points")]:
            vals = filters.get(key, [])
            if vals and len(vals) < 10:
                parts.append(f"{label}: {', '.join(vals[:5])}")
            elif vals:
                parts.append(f"{label}: {len(vals)} seleccionados")
        if parts:
            filter_header = " | ".join(parts)

    # ════════════════════════════════════════════════════════
    # SLIDE 1: TITLE
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_dark_bg(slide)

    # Top accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN_NEON
    line.line.fill.background()

    # System badge
    add_text_box(slide, 0.8, 1.5, 4, 0.8, "SYSTEM JB", 18, GREEN_NEON, True)
    add_text_box(slide, 0.8, 2.3, 8, 1.2, "ENGINE MONITOR", 44, WHITE, True)
    add_text_box(slide, 0.8, 3.4, 8, 0.8, "Analytical Report - Control Administrative", 18, GRAY_TEXT)
    add_text_box(slide, 0.8, 4.2, 8, 0.5, f"{len(records)} registros procesados", 14, GREEN_NEON, False, PP_ALIGN.LEFT)

    if filter_header:
        add_text_box(slide, 0.8, 4.8, 10, 0.6, filter_header, 12, GRAY_TEXT)

    # ════════════════════════════════════════════════════════
    # SLIDE 2: KPI SUMMARY (4 cards)
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)

    add_text_box(slide, 0.5, 0.3, 6, 0.6, "RESUMEN DE INDICADORES", 20, WHITE, True)
    if filter_header:
        add_text_box(slide, 0.5, 0.9, 10, 0.4, filter_header, 10, GRAY_TEXT)

    cards = [
        ("Auditorias Realizadas", f"{total_audits:,}", BLUE_ACCENT),
        ("Fallos Detectados", f"{total_failures:,}", RED_ACCENT),
        ("Tasa de Fracaso", f"{fr_pct:.2f}%", kpi_color(fr_pct)),
        ("Defectos Medicion", f"{total_meas_def:,}", EMERALD_ACCENT),
    ]

    card_w = 2.8
    card_h = 2.2
    gap = 0.3
    total_w = len(cards) * card_w + (len(cards) - 1) * gap
    start_x = (13.333 - total_w) / 2
    y_top = 2.0

    for i, (title, value, color) in enumerate(cards):
        add_card(slide, start_x + i * (card_w + gap), y_top, card_w, card_h, title, value, color)

    # ════════════════════════════════════════════════════════
    # SLIDE 3: CHARTS
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)

    add_text_box(slide, 0.5, 0.3, 6, 0.6, "ANALISIS GRAFICO", 20, WHITE, True)

    if chart_image_path and os.path.exists(chart_image_path):
        slide.shapes.add_picture(
            chart_image_path,
            Inches(1.0), Inches(1.5),
            Inches(11.3), Inches(5.5)
        )
    else:
        add_text_box(slide, 0.5, 3.0, 10, 0.6, "[Grafico no disponible - imagen no proporcionada]",
                     16, GRAY_TEXT, False, PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════
    # SLIDE 4: DETAIL TABLE
    # ════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)

    add_text_box(slide, 0.5, 0.3, 6, 0.6, "MATRIZ DE DETALLE", 20, WHITE, True)
    if filter_header:
        add_text_box(slide, 0.5, 0.8, 10, 0.4, filter_header, 10, GRAY_TEXT)

    # Table: Month | Factory | Buyer | Audits | Failures | FR% | M.Def% | V.Def% | Result
    headers = ["Month", "Factory", "Buyer", "Audits", "Failures", "FR%", "M.Def%", "V.Def%", "Result"]
    col_widths = [1.5, 2.5, 1.8, 1.2, 1.2, 1.2, 1.2, 1.2, 1.0]

    # Aggregate records for display
    from collections import defaultdict
    agg = defaultdict(lambda: {"audits": 0, "failures": 0, "measDef": 0, "measIns": 0, "visDef": 0, "visIns": 0})
    for r in records:
        key = (r.get("mesTexto", ""), r.get("factory", ""), r.get("buyer", ""))
        agg[key]["audits"] += r.get("audits", 0)
        agg[key]["failures"] += r.get("failures", 0)
        agg[key]["measDef"] += r.get("measDef", 0)
        agg[key]["measIns"] += r.get("measIns", 0)
        agg[key]["visDef"] += r.get("visDef", 0)
        agg[key]["visIns"] += r.get("visIns", 0)

    table_data = []
    for (month, factory, buyer), v in sorted(agg.items()):
        fr = v["failures"] / v["audits"] * 100 if v["audits"] > 0 else 0
        mr = v["measDef"] / v["measIns"] * 100 if v["measIns"] > 0 else 0
        vr = v["visDef"] / v["visIns"] * 100 if v["visIns"] > 0 else 0
        result = "PASS" if fr < 10 and mr < 15 and vr < 20 else "FAIL"
        table_data.append([month, factory, buyer, str(v["audits"]), str(v["failures"]),
                          f"{fr:.1f}%", f"{mr:.1f}%", f"{vr:.1f}%", result])

    num_rows = len(table_data) + 1
    num_cols = len(headers)
    tbl_shape = slide.shapes.add_table(num_rows, num_cols,
                                       Inches(0.5), Inches(1.4), Inches(sum(col_widths)), Inches(0.35 * num_rows + 0.3))
    table = tbl_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    # Style header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x00, 0xD0, 0x84)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Style data rows
    for ri, row in enumerate(table_data):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)

            # Alternating row colors
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x22)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1C, 0x21, 0x28)

            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = WHITE if ci < 3 else GRAY_TEXT
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER if ci > 2 else PP_ALIGN.LEFT

            # Result column color
            if ci == 8:
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.bold = True
                    if val == "PASS":
                        p.font.color.rgb = GREEN_NEON
                    else:
                        p.font.color.rgb = RED_ACCENT

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    prs.save(output_path)
    print(f"Report exported: {output_path}")


def generate_sample_data():
    import datetime, random
    rows = []
    factories = ["SAE-A TECHNOTEX, S.A(1)", "SAE-A TECHNOTEX, S.A(2)", "EINS, S.A.(2)"]
    buyers = ["NIKE", "ADIDAS", "PUMA", "UNDER ARMOUR"]
    stages = ["CUTTING", "SEWING", "PACKING"]
    points = ["IN-LINE", "FINAL", "RANDOM"]

    for f in factories:
        for b in buyers:
            for s in stages:
                for p in points:
                    audits = random.randint(50, 500)
                    failures = random.randint(0, int(audits * 0.3))
                    meas_insp = random.randint(20, 200)
                    meas_def = random.randint(0, int(meas_insp * 0.2))
                    vis_insp = random.randint(20, 200)
                    vis_def = random.randint(0, int(vis_insp * 0.2))
                    rows.append({
                        "fecha": datetime.date(2025, random.randint(1, 12), random.randint(1, 28)).isoformat(),
                        "factory": f, "buyer": b, "stage": s, "point": p,
                        "style": "STY-" + str(random.randint(1000, 9999)),
                        "po": "PO-" + str(random.randint(10000, 99999)),
                        "colL": "FABRIC", "colM": "",
                        "audits": audits, "failures": failures,
                        "measIns": meas_insp, "measDef": meas_def,
                        "visIns": vis_insp, "visDef": vis_def,
                        "failureRate": failures/audits, "measDefRate": meas_def/meas_insp, "visDefRate": vis_def/vis_insp,
                        "result": "PASS" if random.random() > 0.3 else "FAIL",
                        "mesFiltroKey": "", "mesTexto": "",
                    })
    return rows


def main():
    parser = argparse.ArgumentParser(description="Generate Engine Monitor PowerPoint report")
    parser.add_argument("input", nargs="?", help="JSON input file (omit for stdin)")
    parser.add_argument("chart", nargs="?", help="Chart image file (PNG)")
    parser.add_argument("-o", "--output", default="engine_monitor.pptx", help="Output .pptx path")
    parser.add_argument("--sample", action="store_true", help="Use sample data")
    parser.add_argument("--filters", help="JSON string with active filters")
    args = parser.parse_args()

    chart_img = args.chart or ""
    filters = None
    if args.filters:
        try:
            filters = json.loads(args.filters)
        except json.JSONDecodeError:
            pass

    if args.sample:
        records = generate_sample_data()
        print("[INFO] Using sample data")
    elif args.input:
        with open(args.input, "r", encoding="utf-8") as f:
            data = json.load(f)
        records = data if isinstance(data, list) else data.get("entries", data.get("data", []))
        print(f"[INFO] Loaded {len(records)} records from {args.input}")
    else:
        data = json.load(sys.stdin)
        records = data if isinstance(data, list) else data.get("entries", data.get("data", []))
        print(f"[INFO] Loaded {len(records)} records from stdin")

    build_presentation(records, chart_img, args.output, filters)


if __name__ == "__main__":
    main()
